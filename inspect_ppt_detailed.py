import pptx

prs = pptx.Presentation('SIH2026-IDEA-Presentation-Format.pptx')
print(f'Slide Width: {prs.slide_width / 914400:.2f} inches, Slide Height: {prs.slide_height / 914400:.2f} inches')

s1 = prs.slides[0]
for idx, shape in enumerate(s1.shapes):
    print(f'Shape {idx}: id={shape.shape_id}, name="{shape.name}", type={shape.shape_type}, pos=({shape.left/914400:.2f}in, {shape.top/914400:.2f}in), size=({shape.width/914400:.2f}in x {shape.height/914400:.2f}in)')
    if shape.has_text_frame:
        for p_i, p in enumerate(shape.text_frame.paragraphs):
            print(f'  P{p_i}: text="{p.text}"')

for i, slide in enumerate(prs.slides):
    print(f'\n=======================================================')
    print(f'=== SLIDE {i+1} ===')
    for s_idx, shape in enumerate(slide.shapes):
        print(f'  Shape {s_idx}: id={shape.shape_id}, name="{shape.name}", type={shape.shape_type}, pos=({shape.left/914400:.2f}in, {shape.top/914400:.2f}in), size=({shape.width/914400:.2f}in x {shape.height/914400:.2f}in)')
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                txt = p.text.strip().replace('\n', ' ')
                if txt:
                    font_size = f"{p.font.size.pt}pt" if p.font and p.font.size else "default"
                    print(f'    [text] ({font_size}): {txt[:100]}')
