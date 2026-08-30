import pptx

prs = pptx.Presentation("SIH2026-IDEA-Presentation-Format.pptx")
print(f"TOTAL SLIDES: {len(prs.slides)}")

for idx, slide in enumerate(prs.slides):
    title = slide.shapes.title.text if slide.shapes.title else "No Title"
    print(f"\n=======================================================")
    print(f"SLIDE {idx+1}: {title}")
    print(f"=======================================================")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                txt = p.text.strip()
                if txt:
                    print(f"  • {txt}")
        if shape.has_table:
            print("  [TABLE]")
            for row in shape.table.rows:
                row_str = " | ".join(cell.text.strip().replace('\n', ' ') for cell in row.cells)
                print(f"    | {row_str} |")
