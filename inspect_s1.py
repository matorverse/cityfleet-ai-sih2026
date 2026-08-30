import pptx

prs = pptx.Presentation('SIH2026_CityFleet_Submission.pptx')
s1 = prs.slides[0]
print(f"Slide 1 Shape count: {len(s1.shapes)}")
for i, s in enumerate(s1.shapes):
    print(f"Shape {i}: name='{s.name}', type={s.shape_type}, left={s.left}, top={s.top}, w={s.width}, h={s.height}")
    if s.has_text_frame:
        for p in s.text_frame.paragraphs:
            print(f"   text: '{p.text}'")
